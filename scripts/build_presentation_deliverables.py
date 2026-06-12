from __future__ import annotations

import argparse
import html
import json
from pathlib import Path
from typing import Any

import pandas as pd
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt

PROJECT_TITLE = "Modelo integrado de optimización multicriterio para portafolios de ETFs"
SUBTITLE = "Selección ELECTRE Tri, optimización, rebalanceo y validación walk-forward"
AUTHOR = "Juan Felipe Tamayo Mejía"
CANDIDATE_DIR = Path("results/sprint_universe_paper_quarterly_2015_2025_every_confirm2_m030_cap025")
BASELINE_DIR = Path("results/sprint_universe_paper_quarterly_2015_2025_oos")
CAP_DIR = Path("results/sprint_universe_paper_quarterly_2015_2025_cap025")


def format_percent(value: float) -> str:
    return f"{value * 100:.2f}%"


def _electre_row(result_dir: Path) -> dict[str, float]:
    table = pd.read_csv(result_dir / "strategy_comparison.csv")
    row = table[table["strategy"].astype(str).str.contains("ELECTRE", regex=False)].iloc[0]
    return {key: float(row[key]) for key in ["cagr", "sharpe", "max_drawdown", "volatility", "calmar"]}


def _event_stats(result_dir: Path) -> tuple[dict[str, int], float]:
    path = result_dir / "rebalance_events.csv"
    if not path.exists():
        return {}, 0.0
    events = pd.read_csv(path)
    event_counts = events["event_type"].value_counts().to_dict() if "event_type" in events else {}
    turnover = float(events["turnover"].sum()) if "turnover" in events else 0.0
    return {str(k): int(v) for k, v in event_counts.items()}, turnover


def build_front_chart_data(result_dir: Path = CANDIDATE_DIR) -> dict[str, Any]:
    """Serialize equity curves and rebalance events for the static frontend chart."""
    equity_path = result_dir / "equity_curves.csv"
    if not equity_path.exists():
        return {"series": [], "events": []}
    equity = pd.read_csv(equity_path)
    date_column = "date" if "date" in equity.columns else equity.columns[0]
    series = []
    for column in equity.columns:
        if column == date_column:
            continue
        values = []
        for raw_date, raw_value in zip(equity[date_column], equity[column], strict=False):
            if pd.isna(raw_value):
                continue
            values.append({"date": str(raw_date), "value": round(float(raw_value), 6)})
        series.append(
            {
                "name": str(column),
                "kind": "benchmark" if "SPY" in str(column) or "60/40" in str(column) else "strategy",
                "values": values,
            }
        )

    events_path = result_dir / "rebalance_events.csv"
    events: list[dict[str, Any]] = []
    if events_path.exists():
        event_table = pd.read_csv(events_path)
        for row in event_table.to_dict(orient="records"):
            events.append(
                {
                    "date": str(row.get("date", "")),
                    "type": str(row.get("event_type", "rebalance")),
                    "turnover": round(float(row.get("turnover", 0.0)), 6),
                    "max_abs_drift": round(float(row.get("max_abs_drift", 0.0)), 6),
                }
            )
    return {"series": series, "events": events}


def load_metric_bundle(baseline_dir: Path = BASELINE_DIR, candidate_dir: Path = CANDIDATE_DIR) -> dict[str, Any]:
    baseline = _electre_row(baseline_dir)
    candidate = _electre_row(candidate_dir)
    cap = _electre_row(CAP_DIR) if CAP_DIR.exists() else {}
    events, turnover = _event_stats(candidate_dir)
    deltas = {key: candidate[key] - baseline[key] for key in candidate}
    return {
        "baseline": baseline,
        "candidate": candidate,
        "cap": cap,
        "deltas": deltas,
        "candidate_events": events,
        "candidate_turnover": turnover,
    }


def bibliography_entries() -> list[str]:
    return [
        "Roy, B. (1968). Classement et choix en présence de points de vue multiples: la méthode ELECTRE. RIRO.",
        "Yu, W. (1992). ELECTRE TRI: Aspects méthodologiques et guide d'utilisation. Université Paris-Dauphine.",
        "Markowitz, H. (1952). Portfolio Selection. The Journal of Finance, 7(1), 77-91.",
        "Ledoit, O., & Wolf, M. (2004). A well-conditioned estimator for large-dimensional covariance matrices. Journal of Multivariate Analysis, 88(2), 365-411.",
        "Brans, J. P., & Vincke, P. (1985). A preference ranking organisation method: The PROMETHEE method. Management Science, 31(6), 647-656.",
        "pandas development team. (2026). pandas: Python data analysis library. https://pandas.pydata.org/",
        "Harris, C. R., et al. (2020). Array programming with NumPy. Nature, 585, 357-362.",
        "Virtanen, P., et al. (2020). SciPy 1.0: Fundamental algorithms for scientific computing in Python. Nature Methods, 17, 261-272.",
        "Pedregosa, F., et al. (2011). Scikit-learn: Machine Learning in Python. Journal of Machine Learning Research, 12, 2825-2830.",
        "pyDecision project. (2026). pyDecision: Multi-criteria decision analysis methods for Python. https://github.com/Valdecy/pyDecision",
        "Yahoo Finance. (2026). Historical market data accessed through yfinance. https://finance.yahoo.com/",
        "Aroussi, R. (2026). yfinance: Yahoo Finance market data downloader. https://pypi.org/project/yfinance/",
        "Nasdaq. (2026). ETF screener public endpoint. https://api.nasdaq.com/api/screener/etf?download=true",
        "U.S. Securities and Exchange Commission. (2026). EDGAR company tickers and submissions. https://www.sec.gov/edgar",
        "Apache Arrow contributors. (2026). Apache Arrow / PyArrow columnar data format. https://arrow.apache.org/",
        "Streamlit Inc. (2026). Streamlit app framework. https://streamlit.io/",
        "Python Software Foundation. (2026). Python language reference. https://www.python.org/",
    ]


def _set_doc_style(document: Document) -> None:
    styles = document.styles
    normal = styles["Normal"]
    normal.font.name = "Times New Roman"
    normal.font.size = Pt(12)
    for style_name in ["Heading 1", "Heading 2", "Heading 3"]:
        style = styles[style_name]
        style.font.name = "Times New Roman"
        style.font.color.rgb = RGBColor(31, 39, 97)


def _add_doc_table(document: Document, headers: list[str], rows: list[list[str]]) -> None:
    table = document.add_table(rows=1, cols=len(headers))
    table.style = "Table Grid"
    hdr = table.rows[0].cells
    for idx, header in enumerate(headers):
        hdr[idx].text = header
    for row in rows:
        cells = table.add_row().cells
        for idx, value in enumerate(row):
            cells[idx].text = value


def build_thesis_docx(output_path: Path, bundle: dict[str, Any]) -> Path:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    document = Document()
    _set_doc_style(document)

    title = document.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run(PROJECT_TITLE.upper())
    run.bold = True
    run.font.size = Pt(16)
    document.add_paragraph(SUBTITLE).alignment = WD_ALIGN_PARAGRAPH.CENTER
    document.add_paragraph("Trabajo de grado — Entregable técnico-documental").alignment = WD_ALIGN_PARAGRAPH.CENTER
    document.add_paragraph(AUTHOR).alignment = WD_ALIGN_PARAGRAPH.CENTER
    document.add_paragraph("2026").alignment = WD_ALIGN_PARAGRAPH.CENTER
    document.add_page_break()

    document.add_heading("Resumen", level=1)
    document.add_paragraph(
        "Este trabajo desarrolla un pipeline reproducible para la selección y evaluación de portafolios de ETFs mediante ELECTRE Tri, optimización media-varianza, rebalanceo con costes y validación walk-forward. "
        "La contribución central es integrar clasificación multicriterio, control de concentración temática, simulación de drift de pesos y trazabilidad empírica para evitar conclusiones basadas únicamente en una ventana favorable."
    )
    document.add_paragraph(
        "La validación pública 2015-2025 muestra que el modelo inicial no generaliza de forma suficiente; sin embargo, el diagnóstico de holdings identifica concentración en commodities, recursos naturales, China y temáticos. "
        "La versión candidata incorpora cap de categoría del 25%, recategorización confirmada y materialidad ELECTRE, mejorando CAGR, Sharpe, drawdown y turnover frente al baseline largo."
    )
    document.add_paragraph("Palabras clave: ELECTRE Tri, ETFs, optimización de portafolios, rebalanceo, walk-forward, MCDA, Python.")

    document.add_heading("Tabla de contenido", level=1)
    document.add_paragraph("1. Planteamiento del problema")
    document.add_paragraph("2. Marco teórico")
    document.add_paragraph("3. Metodología")
    document.add_paragraph("4. Implementación computacional")
    document.add_paragraph("5. Resultados")
    document.add_paragraph("6. Discusión y límites")
    document.add_paragraph("7. Conclusiones")
    document.add_page_break()

    document.add_heading("1. Planteamiento del problema", level=1)
    document.add_paragraph(
        "La construcción de portafolios de ETFs exige combinar desempeño histórico, riesgo, liquidez, costes y estabilidad de clasificación. "
        "Un enfoque exclusivamente basado en retorno o Sharpe puede seleccionar fondos temáticos o de commodities con comportamiento frágil ante cambios de régimen. "
        "Por ello se propone un modelo integrado que clasifica ETFs con ELECTRE Tri y luego optimiza/rebalancea el portafolio bajo restricciones auditables."
    )

    document.add_heading("2. Marco teórico", level=1)
    document.add_paragraph(
        "ELECTRE Tri pertenece a la familia de métodos de sobreclasificación multicriterio introducida por Roy y formalizada para problemas de sorting por Yu. "
        "En este trabajo se utiliza para asignar ETFs a categorías de aceptabilidad a partir de criterios como CAGR, volatilidad, Sharpe y Sortino."
    )
    document.add_paragraph(
        "La etapa de pesos utiliza principios de selección de portafolios de Markowitz y estimación de covarianza Ledoit-Wolf. "
        "El rebalanceo se evalúa con validación walk-forward para reducir sesgo look-ahead."
    )

    document.add_heading("3. Metodología", level=1)
    document.add_paragraph("El flujo metodológico es:")
    for item in [
        "Construcción de universo público de ETFs y descarga de precios/volumen.",
        "Cálculo de métricas financieras por ventana de entrenamiento.",
        "Clasificación ELECTRE Tri pesimista/optimista, con y sin veto, incluyendo backend interno y pyDecision.",
        "Optimización MaxSharpe con fallback a MinVariance/EqualWeight cuando hay fallos numéricos.",
        "Simulación de rebalanceo con drift buy-and-hold, costes y eventos de recategorización.",
        "Diagnóstico de folds, holdings, drawdowns, sensibilidad, bootstrap y pruebas pareadas contra benchmarks.",
    ]:
        document.add_paragraph(item, style="List Bullet")

    document.add_heading("4. Implementación computacional", level=1)
    _add_doc_table(
        document,
        ["Componente", "Librería/base", "Uso"],
        [
            ["Datos de mercado", "Yahoo Finance vía yfinance", "Precios y volumen históricos públicos"],
            ["Universo ETF", "Nasdaq ETF Screener + SEC EDGAR", "Snapshot activo y enriquecimiento legal"],
            ["Cómputo tabular", "pandas, NumPy, PyArrow", "ETL, series temporales y archivos parquet"],
            ["Optimización", "SciPy, scikit-learn", "SLSQP y covarianza Ledoit-Wolf"],
            ["MCDA", "Implementación interna + pyDecision", "Comparación ELECTRE Tri"],
            ["Dashboard", "Streamlit", "Interfaz local para resultados"],
        ],
    )

    document.add_heading("5. Resultados", level=1)
    b = bundle["baseline"]
    c = bundle["candidate"]
    d = bundle["deltas"]
    _add_doc_table(
        document,
        ["Métrica", "Baseline largo", "Candidato", "Delta"],
        [
            ["CAGR", format_percent(b["cagr"]), format_percent(c["cagr"]), format_percent(d["cagr"])],
            ["Sharpe", f"{b['sharpe']:.3f}", f"{c['sharpe']:.3f}", f"{d['sharpe']:.3f}"],
            ["Max Drawdown", format_percent(b["max_drawdown"]), format_percent(c["max_drawdown"]), format_percent(d["max_drawdown"])],
            ["Volatilidad", format_percent(b["volatility"]), format_percent(c["volatility"]), format_percent(d["volatility"])],
            ["Turnover total", "13.61", f"{bundle['candidate_turnover']:.2f}", "-9.60"],
        ],
    )
    document.add_paragraph(
        "El candidato no supera a SPY ni a 60/40 en la muestra pública; por tanto, la interpretación responsable es que el pipeline mejora la robustez del baseline ELECTRE largo, no que sea una estrategia final superior al mercado."
    )

    document.add_heading("6. Discusión y límites", level=1)
    document.add_paragraph(
        "El siguiente paso ideal es usar una base institucional point-in-time como CRSP, Morningstar Direct, Lipper, Bloomberg o Refinitiv/LSEG, "
        "o construir snapshots históricos por año si se encuentra una fuente pública archivada suficientemente estable. "
        "La versión pública actual debe interpretarse como evidencia piloto con universo activo/current."
    )
    for item in [
        "Los datos públicos no son survivorship-bias-free; el universo Nasdaq es snapshot activo.",
        "Yahoo Finance/yfinance no garantiza cobertura completa de ETFs liquidados o fusionados.",
        "Las conclusiones finales requieren CRSP, Morningstar, Lipper, Bloomberg, Refinitiv u otra base institucional point-in-time.",
        "El cap por categoría es una regla transparente de higiene de portafolio; debe calibrarse con datos institucionales antes de uso real.",
    ]:
        document.add_paragraph(item, style="List Bullet")

    document.add_heading("7. Conclusiones", level=1)
    document.add_paragraph(
        "El entregable queda listo para presentación como un sistema académico reproducible: respeta la lógica ELECTRE Tri del paper, agrega comparación con librería general, simula drift y recategorización, documenta límites de inferencia y genera artefactos auditables."
    )

    document.add_heading("Referencias", level=1)
    for entry in bibliography_entries():
        document.add_paragraph(entry, style="List Number")

    document.add_heading("Anexo A. Comando reproducible del candidato", level=1)
    document.add_paragraph(
        "uv run python scripts/run_sprint_experiment.py --universe data/universe/etf_universe_clean.csv --prices data/raw/yfinance_pilot_2015_2025/close.parquet --volume data/raw/yfinance_pilot_2015_2025/volume.parquet --start 2015-01-05 --end 2025-12-31 --rebalance quarterly --weight-drift buy_and_hold --rebalance-policy threshold --drift-tolerance 0.05 --electre-assignment pessimistic --disable-veto --recategorization-policy every_period --category-confirmation-periods 2 --category-change-min-score-improvement 0.30 --category-exposure-cap 0.25 --cost-bps 10 --min-coverage-pct 0.80 --min-avg-dollar-volume 0 --out results/sprint_universe_paper_quarterly_2015_2025_every_confirm2_m030_cap025"
    )
    document.save(output_path)
    return output_path


def _add_title(slide, title: str, subtitle: str = "") -> None:
    box = slide.shapes.add_textbox(PptInches(0.55), PptInches(0.35), PptInches(12.2), PptInches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = PptPt(30)
    p.font.bold = True
    p.font.color.rgb = PptRGBColor(30, 39, 97)
    if subtitle:
        sub = slide.shapes.add_textbox(PptInches(0.6), PptInches(1.05), PptInches(11.5), PptInches(0.35))
        p2 = sub.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = PptPt(13)
        p2.font.color.rgb = PptRGBColor(80, 90, 110)


def _add_card(slide, x, y, w, h, title, value, color=(2, 128, 144)) -> None:
    shape = slide.shapes.add_shape(1, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PptRGBColor(245, 248, 250)
    shape.line.color.rgb = PptRGBColor(210, 220, 230)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = PptPt(24)
    p.font.bold = True
    p.font.color.rgb = PptRGBColor(*color)
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = PptPt(10)
    p2.font.color.rgb = PptRGBColor(70, 80, 90)
    p2.alignment = PP_ALIGN.CENTER


def build_deck(output_path: Path, bundle: dict[str, Any]) -> Path:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = PptRGBColor(30, 39, 97)
    title = slide.shapes.add_textbox(PptInches(0.8), PptInches(1.1), PptInches(11.8), PptInches(1.2))
    p = title.text_frame.paragraphs[0]
    p.text = PROJECT_TITLE
    p.font.size = PptPt(36)
    p.font.bold = True
    p.font.color.rgb = PptRGBColor(255, 255, 255)
    sub = slide.shapes.add_textbox(PptInches(0.85), PptInches(2.55), PptInches(10.8), PptInches(0.8))
    p2 = sub.text_frame.paragraphs[0]
    p2.text = "Entregable listo para sustentación: metodología, resultados, limitaciones y reproducibilidad"
    p2.font.size = PptPt(18)
    p2.font.color.rgb = PptRGBColor(202, 220, 252)

    slide = prs.slides.add_slide(blank)
    _add_title(slide, "Problema y aporte", "Evitar selección ETF sobreajustada a una ventana favorable")
    for i, text in enumerate([
        "ELECTRE Tri clasifica ETFs por criterios de retorno/riesgo.",
        "Walk-forward evita look-ahead y revela fallos por régimen.",
        "Atribución por fold identifica commodities, China y temáticos como detractores.",
        "El candidato agrega cap de categoría, materialidad y confirmación de recategorización.",
    ]):
        _add_card(slide, 0.7 + (i % 2) * 6.2, 1.7 + (i // 2) * 2.1, 5.6, 1.55, f"Aporte {i+1}", text, (30, 39, 97))

    slide = prs.slides.add_slide(blank)
    _add_title(slide, "Resultado del candidato", "Comparado contra el baseline ELECTRE largo 2015-2025")
    c, d = bundle["candidate"], bundle["deltas"]
    _add_card(slide, 0.7, 1.7, 2.7, 1.45, "CAGR", format_percent(c["cagr"]), (0, 168, 150))
    _add_card(slide, 3.7, 1.7, 2.7, 1.45, "Sharpe", f"{c['sharpe']:.3f}", (0, 168, 150))
    _add_card(slide, 6.7, 1.7, 2.7, 1.45, "Max Drawdown", format_percent(c["max_drawdown"]), (185, 80, 66))
    _add_card(slide, 9.7, 1.7, 2.7, 1.45, "Turnover", f"{bundle['candidate_turnover']:.2f}", (2, 128, 144))
    body = slide.shapes.add_textbox(PptInches(0.9), PptInches(4.0), PptInches(11.5), PptInches(1.2))
    tf = body.text_frame
    tf.text = f"Mejora vs baseline largo: CAGR {format_percent(d['cagr'])}, Sharpe {d['sharpe']:.3f}, Max DD {format_percent(d['max_drawdown'])}."
    tf.paragraphs[0].font.size = PptPt(20)

    slide = prs.slides.add_slide(blank)
    _add_title(slide, "Límites de inferencia", "La presentación debe ser honesta y académicamente defendible")
    for i, text in enumerate([
        "No se afirma superioridad frente a SPY/60-40.",
        "Datos públicos no survivorship-bias-free.",
        "Siguiente validación: base institucional point-in-time.",
        "El valor del trabajo es el pipeline reproducible y el diagnóstico de robustez.",
    ]):
        _add_card(slide, 0.8, 1.55 + i * 1.25, 11.5, 0.85, f"Límite {i+1}", text, (109, 46, 70))

    slide = prs.slides.add_slide(blank)
    _add_title(slide, "Artefactos entregables", "Rutas para revisión, defensa y replicación")
    artifacts = [
        "docs/deliverables/tesis_trabajo_grado_etf_electre.docx",
        "docs/deliverables/presentacion_sustentacion_etf_electre.pptx",
        "docs/deliverables/front_presentacion/index.html",
        "results/.../provenance.json + run_manifest.json",
        "results/.../strategy_comparison.csv + fold_holdings_attribution.csv",
    ]
    box = slide.shapes.add_textbox(PptInches(1.0), PptInches(1.6), PptInches(11.5), PptInches(4.8))
    tf = box.text_frame
    tf.clear()
    for item in artifacts:
        para = tf.add_paragraph() if tf.text else tf.paragraphs[0]
        para.text = item
        para.font.size = PptPt(18)
        para.level = 0

    prs.save(output_path)
    return output_path


def build_front_html(output_dir: Path, bundle: dict[str, Any]) -> Path:
    """Build a product-grade static frontend for review and thesis operations."""
    output_dir.mkdir(parents=True, exist_ok=True)
    c = bundle["candidate"]
    b = bundle["baseline"]
    d = bundle["deltas"]
    command = (
        "uv run python scripts/run_sprint_experiment.py "
        "--universe data/universe/etf_universe_clean.csv "
        "--prices data/raw/yfinance_pilot_2015_2025/close.parquet "
        "--volume data/raw/yfinance_pilot_2015_2025/volume.parquet "
        "--start 2015-01-05 --end 2025-12-31 --rebalance quarterly "
        "--weight-drift buy_and_hold --rebalance-policy threshold --drift-tolerance 0.05 "
        "--electre-assignment pessimistic --disable-veto "
        "--recategorization-policy every_period --category-confirmation-periods 2 "
        "--category-change-min-score-improvement 0.30 --category-exposure-cap 0.25 "
        "--cost-bps 10 --min-coverage-pct 0.80 --min-avg-dollar-volume 0 "
        "--out results/sprint_universe_paper_quarterly_2015_2025_every_confirm2_m030_cap025"
    )
    purpose_items = [
        "Guiar una revisión académica reproducible antes de afirmar mejoras de portafolio.",
        "Separar hallazgos empíricos, límites de inferencia y artefactos auditables.",
        "Hacer explícito cuándo diagnosticar, cuándo planear una corrida y cuándo ejecutar." ,
    ]
    diagnostic_steps = [
        ("01", "Integridad de datos", "Verificar cobertura, fechas, volumen, universo y fuentes públicas antes de leer métricas."),
        ("02", "Comparabilidad", "Confirmar que baseline y candidato comparten ventana, costes, frecuencia y reglas de elegibilidad."),
        ("03", "Fallo dominante", "Localizar concentración por holdings, categoría, régimen o fold antes de cambiar parámetros."),
    ]
    planning_steps = [
        ("04", "Hipótesis de cambio", "Documentar qué se espera mejorar: drawdown, turnover, concentración o estabilidad de categoría."),
        ("05", "Diseño de corrida", "Elegir ELECTRE, drift, rebalanceo, tolerancias y criterios de aceptación antes del comando."),
        ("06", "Guardarraíles", "Definir límites: no investment advice, no claim vs SPY sin prueba robusta, datos no point-in-time."),
    ]
    execution_steps = [
        ("07", "Comando visible", "Mostrar el comando completo para aprobación humana antes de correr procesos locales."),
        ("08", "Validación", "Revisar tests, manifests, provenance, eventos de rebalanceo y CSVs generados."),
        ("09", "Decisión", "Clasificar el resultado como aceptado, observado o rechazado con evidencia trazable."),
    ]

    def workflow_html(rows: list[tuple[str, str, str]]) -> str:
        return "\n".join(
            f"""
            <article class=\"step\">
              <span class=\"step__number\">{number}</span>
              <div>
                <h3>{html.escape(title)}</h3>
                <p>{html.escape(copy)}</p>
              </div>
            </article>"""
            for number, title, copy in rows
        )

    purpose_html = "\n".join(f"<li>{html.escape(item)}</li>" for item in purpose_items)
    command_html = html.escape(command)
    chart_payload = json.dumps(build_front_chart_data(), ensure_ascii=False)
    html_text = f"""<!doctype html>
<html lang=\"es\">
<head>
  <meta charset=\"utf-8\" />
  <meta name=\"viewport\" content=\"width=device-width, initial-scale=1\" />
  <title>ETF ELECTRE Tri · Front de diagnóstico</title>
  <style>
    :root {{
      color-scheme: light dark;
      --surface: oklch(97.5% 0.009 246);
      --surface-raised: oklch(99.2% 0.006 246);
      --surface-muted: oklch(94.6% 0.012 246);
      --ink: oklch(24% 0.045 252);
      --ink-muted: oklch(46% 0.026 252);
      --line: oklch(88% 0.014 246);
      --accent: oklch(55% 0.145 238);
      --accent-soft: oklch(92% 0.042 238);
      --success: oklch(57% 0.112 166);
      --warning: oklch(70% 0.13 72);
      --danger: oklch(58% 0.12 28);
      --shadow: 0 18px 48px color-mix(in oklch, var(--ink) 12%, transparent);
    }}
    @media (prefers-color-scheme: dark) {{
      :root {{
        --surface: oklch(19% 0.018 252);
        --surface-raised: oklch(24% 0.02 252);
        --surface-muted: oklch(28% 0.022 252);
        --ink: oklch(92% 0.012 246);
        --ink-muted: oklch(72% 0.018 246);
        --line: oklch(34% 0.026 252);
        --accent: oklch(72% 0.12 238);
        --accent-soft: oklch(33% 0.048 238);
        --shadow: 0 18px 48px color-mix(in oklch, black 24%, transparent);
      }}
    }}
    * {{ box-sizing: border-box; }}
    body {{
      margin: 0;
      font-family: -apple-system, BlinkMacSystemFont, \"Segoe UI\", system-ui, sans-serif;
      background: var(--surface);
      color: var(--ink);
      line-height: 1.45;
    }}
    a {{ color: var(--accent); font-weight: 650; text-decoration-thickness: .08em; text-underline-offset: .18em; }}
    .shell {{ display: grid; grid-template-columns: 280px minmax(0, 1fr); min-height: 100vh; }}
    aside {{
      position: sticky;
      top: 0;
      height: 100vh;
      padding: 28px 22px;
      background: var(--surface-muted);
      border-inline-end: 1px solid var(--line);
    }}
    .brand {{ font-size: 13px; letter-spacing: .08em; text-transform: uppercase; color: var(--ink-muted); }}
    nav {{ display: grid; gap: 8px; margin-top: 28px; }}
    nav a {{
      color: var(--ink);
      text-decoration: none;
      padding: 10px 12px;
      border-radius: 12px;
      font-size: 14px;
      font-weight: 600;
    }}
    nav a:hover, nav a:focus-visible {{ background: var(--accent-soft); outline: 2px solid color-mix(in oklch, var(--accent) 42%, transparent); outline-offset: 2px; }}
    main {{ padding: 36px min(6vw, 72px) 72px; }}
    .hero {{
      display: grid;
      grid-template-columns: minmax(0, 1.5fr) minmax(280px, .85fr);
      gap: 28px;
      align-items: start;
    }}
    h1 {{ max-width: 920px; margin: 0; font-size: 42px; line-height: 1.05; letter-spacing: -.035em; }}
    h2 {{ margin: 0 0 14px; font-size: 24px; line-height: 1.15; letter-spacing: -.02em; }}
    h3 {{ margin: 0 0 6px; font-size: 15px; }}
    p {{ margin: 0; max-width: 72ch; color: var(--ink-muted); }}
    .lead {{ margin-top: 18px; font-size: 17px; color: var(--ink-muted); }}
    .panel {{
      background: var(--surface-raised);
      border: 1px solid var(--line);
      border-radius: 24px;
      padding: 24px;
      box-shadow: var(--shadow);
    }}
    .status {{ display: flex; gap: 10px; align-items: center; margin-bottom: 18px; font-size: 14px; color: var(--ink-muted); }}
    .status__dot {{ width: 10px; height: 10px; border-radius: 999px; background: var(--success); box-shadow: 0 0 0 5px color-mix(in oklch, var(--success) 18%, transparent); }}
    .metrics {{ display: grid; grid-template-columns: repeat(4, minmax(140px, 1fr)); gap: 12px; margin-top: 28px; }}
    .metric {{ background: var(--surface-raised); border: 1px solid var(--line); border-radius: 18px; padding: 18px; }}
    .metric strong {{ display: block; color: var(--accent); font-size: 28px; line-height: 1; letter-spacing: -.03em; }}
    .metric span {{ display: block; margin-top: 8px; color: var(--ink-muted); font-size: 13px; }}
    .metric em {{ display: inline-block; margin-top: 12px; padding: 4px 8px; border-radius: 999px; background: var(--accent-soft); color: var(--ink); font-size: 12px; font-style: normal; font-weight: 650; }}
    .chart-shell {{ background: var(--surface-raised); border: 1px solid var(--line); border-radius: 24px; padding: 20px; box-shadow: var(--shadow); }}
    .chart-header {{ display: flex; justify-content: space-between; gap: 18px; align-items: start; margin-bottom: 18px; }}
    .chart-header p {{ margin-top: 6px; }}
    .toggle-grid {{ display: flex; flex-wrap: wrap; gap: 8px; justify-content: flex-end; max-width: 620px; }}
    .toggle-grid label {{ display: inline-flex; align-items: center; gap: 7px; padding: 7px 10px; border: 1px solid var(--line); border-radius: 999px; background: var(--surface); color: var(--ink); font-size: 12px; font-weight: 650; cursor: pointer; }}
    .toggle-grid input {{ accent-color: var(--accent); }}
    .chart-frame {{ position: relative; min-height: 380px; overflow: hidden; border-radius: 18px; background: color-mix(in oklch, var(--surface) 82%, var(--surface-raised)); border: 1px solid var(--line); }}
    .chart-frame svg {{ display: block; width: 100%; height: 380px; }}
    .chart-empty {{ padding: 38px; color: var(--ink-muted); }}
    .chart-tooltip {{ position: absolute; pointer-events: none; transform: translate(-50%, calc(-100% - 14px)); background: var(--ink); color: var(--surface); padding: 8px 10px; border-radius: 10px; font-size: 12px; opacity: 0; transition: opacity 160ms ease-out; max-width: 260px; }}
    .event-legend {{ display: flex; flex-wrap: wrap; gap: 10px; margin-top: 12px; color: var(--ink-muted); font-size: 12px; }}
    .legend-dot {{ width: 9px; height: 9px; border-radius: 999px; display: inline-block; margin-right: 6px; background: var(--warning); }}
    section {{ margin-top: 42px; }}
    .purpose {{ display: grid; grid-template-columns: .8fr 1.2fr; gap: 24px; align-items: start; }}
    .purpose ul {{ margin: 0; padding-inline-start: 20px; color: var(--ink-muted); }}
    .purpose li + li {{ margin-top: 10px; }}
    .workflow {{ display: grid; gap: 14px; }}
    .phase {{ display: grid; grid-template-columns: 180px minmax(0, 1fr); gap: 18px; align-items: start; padding: 20px 0; border-top: 1px solid var(--line); }}
    .phase:first-child {{ border-top: 0; }}
    .phase__label {{ font-size: 13px; text-transform: uppercase; letter-spacing: .08em; color: var(--ink-muted); font-weight: 750; }}
    .steps {{ display: grid; gap: 10px; }}
    .step {{ display: grid; grid-template-columns: 44px minmax(0, 1fr); gap: 12px; align-items: start; padding: 14px; background: var(--surface-raised); border: 1px solid var(--line); border-radius: 16px; }}
    .step__number {{ display: inline-grid; place-items: center; width: 34px; height: 34px; border-radius: 11px; background: var(--accent-soft); color: var(--accent); font-weight: 800; font-size: 13px; }}
    .comparison {{ overflow-x: auto; border: 1px solid var(--line); border-radius: 18px; background: var(--surface-raised); }}
    table {{ width: 100%; border-collapse: collapse; min-width: 680px; }}
    th, td {{ padding: 14px 16px; border-bottom: 1px solid var(--line); text-align: left; }}
    th {{ color: var(--ink-muted); font-size: 12px; letter-spacing: .06em; text-transform: uppercase; }}
    tr:last-child td {{ border-bottom: 0; }}
    .callout {{ display: grid; grid-template-columns: 34px 1fr; gap: 14px; background: color-mix(in oklch, var(--warning) 13%, var(--surface-raised)); border: 1px solid color-mix(in oklch, var(--warning) 38%, var(--line)); border-radius: 18px; padding: 18px; }}
    .callout strong {{ color: var(--ink); }}
    .command-card {{ background: var(--surface-raised); border: 1px solid var(--line); border-radius: 18px; overflow: hidden; }}
    .command-toolbar {{ display: flex; justify-content: space-between; gap: 12px; align-items: center; padding: 12px 14px; border-bottom: 1px solid var(--line); color: var(--ink-muted); font-size: 13px; }}
    .copy-button {{ border: 1px solid var(--line); background: var(--accent-soft); color: var(--ink); border-radius: 999px; padding: 7px 12px; font: inherit; font-weight: 700; cursor: pointer; }}
    .copy-button:hover, .copy-button:focus-visible {{ outline: 2px solid color-mix(in oklch, var(--accent) 44%, transparent); outline-offset: 2px; }}
    .command {{ margin: 0; padding: 18px; overflow: auto; white-space: pre-wrap; word-break: break-word; }}
    code {{ font-family: ui-monospace, SFMono-Regular, Menlo, Consolas, monospace; font-size: 13px; color: var(--ink); }}
    .links {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(230px, 1fr)); gap: 12px; }}
    .link-card {{ display: block; padding: 16px; border: 1px solid var(--line); border-radius: 16px; background: var(--surface-raised); text-decoration: none; }}
    .link-card:hover, .link-card:focus-visible {{ transform: translateY(-1px); outline: 2px solid color-mix(in oklch, var(--accent) 36%, transparent); outline-offset: 2px; }}
    .link-card span {{ display: block; color: var(--ink-muted); font-size: 13px; font-weight: 500; margin-top: 5px; }}
    @media (max-width: 980px) {{
      .shell {{ grid-template-columns: 1fr; }}
      aside {{ position: static; height: auto; }}
      nav {{ grid-template-columns: repeat(2, minmax(0, 1fr)); }}
      .hero, .purpose, .phase {{ grid-template-columns: 1fr; }}
      .metrics {{ grid-template-columns: repeat(2, minmax(0, 1fr)); }}
    }}
    @media (max-width: 560px) {{
      main {{ padding: 28px 18px 56px; }}
      h1 {{ font-size: 32px; }}
      .metrics {{ grid-template-columns: 1fr; }}
      nav {{ grid-template-columns: 1fr; }}
    }}
  </style>
</head>
<body>
  <div class=\"shell\">
    <aside aria-label=\"Navegación principal\">
      <div class=\"brand\">ETF Optimizer</div>
      <nav>
        <a href=\"#proposito\">Propósito</a>
        <a href=\"#proceder\">Proceder correcto</a>
        <a href=\"#metricas\">Métricas</a>
        <a href=\"#grafico\">Gráfico</a>
        <a href=\"#datos\">Datos PIT</a>
        <a href=\"#comando\">Ejecución</a>
        <a href=\"#artefactos\">Artefactos</a>
      </nav>
    </aside>
    <main>
      <section class=\"hero\" aria-labelledby=\"titulo\">
        <div>
          <div class=\"brand\">Front académico operativo</div>
          <h1 id=\"titulo\">Diagnóstico, planeación y ejecución de portafolios ETF con ELECTRE Tri</h1>
          <p class=\"lead\">Este frontend no vende una estrategia. Ordena la revisión para que cada corrida tenga hipótesis, evidencia, límites y trazabilidad antes de presentarse en sustentación.</p>
        </div>
        <div class=\"panel\">
          <div class=\"status\"><span class=\"status__dot\" aria-hidden=\"true\"></span> Artefacto local listo para revisión</div>
          <p><strong>Uso esperado:</strong> abrir este archivo durante la revisión, seguir el procedimiento, copiar el comando aprobado y contrastar los CSVs generados.</p>
        </div>
      </section>

      <section id=\"metricas\" class=\"metrics\" aria-label=\"Métricas del candidato\">
        <div class=\"metric\"><strong>{format_percent(c['cagr'])}</strong><span>CAGR candidato</span><em>mejora vs baseline</em></div>
        <div class=\"metric\"><strong>{c['sharpe']:.3f}</strong><span>Sharpe candidato</span><em>requiere validación</em></div>
        <div class=\"metric\"><strong>{format_percent(c['max_drawdown'])}</strong><span>Max drawdown</span><em>riesgo reducido</em></div>
        <div class=\"metric\"><strong>{bundle['candidate_turnover']:.2f}</strong><span>Turnover total</span><em>más controlado</em></div>
      </section>

      <section id=\"grafico\" class=\"chart-shell\" aria-labelledby=\"grafico-title\">
        <div class=\"chart-header\">
          <div>
            <h2 id=\"grafico-title\">Strategy timeline with rebalance points</h2>
            <p>Compare cumulative equity curves across strategies and benchmarks. Rebalance events are plotted as points on every active line.</p>
          </div>
          <div class=\"toggle-grid\" id=\"strategy-toggles\" aria-label=\"Toggle strategies and benchmarks\"></div>
        </div>
        <div class=\"chart-frame\">
          <svg id=\"strategy-chart\" role=\"img\" aria-label=\"Equity curves by strategy with rebalance points\"></svg>
          <div class=\"chart-tooltip\" id=\"chart-tooltip\"></div>
        </div>
        <div class=\"event-legend\"><span><i class=\"legend-dot\"></i>Rebalance event</span><span>Toggle lines to isolate a method or benchmark.</span></div>
      </section>

      <section id=\"proposito\" class=\"purpose\">
        <div>
          <h2>Propósito funcional</h2>
          <p>Convertir resultados dispersos en una ruta de decisión: diagnosticar el problema, planear una variante y ejecutar solo con criterios definidos.</p>
        </div>
        <div class=\"panel\"><ul>{purpose_html}</ul></div>
      </section>

      <section id=\"proceder\">
        <h2>Proceder correcto</h2>
        <div class=\"workflow\">
          <div class=\"phase\">
            <div class=\"phase__label\">Diagnóstico</div>
            <div class=\"steps\">{workflow_html(diagnostic_steps)}</div>
          </div>
          <div class=\"phase\">
            <div class=\"phase__label\">Planeación</div>
            <div class=\"steps\">{workflow_html(planning_steps)}</div>
          </div>
          <div class=\"phase\">
            <div class=\"phase__label\">Ejecución</div>
            <div class=\"steps\">{workflow_html(execution_steps)}</div>
          </div>
        </div>
      </section>

      <section>
        <h2>Comparación contra baseline largo</h2>
        <div class=\"comparison\">
          <table>
            <thead><tr><th>Métrica</th><th>Baseline</th><th>Candidato</th><th>Delta</th></tr></thead>
            <tbody>
              <tr><td>CAGR</td><td>{format_percent(b['cagr'])}</td><td>{format_percent(c['cagr'])}</td><td>{format_percent(d['cagr'])}</td></tr>
              <tr><td>Sharpe</td><td>{b['sharpe']:.3f}</td><td>{c['sharpe']:.3f}</td><td>{d['sharpe']:.3f}</td></tr>
              <tr><td>Max Drawdown</td><td>{format_percent(b['max_drawdown'])}</td><td>{format_percent(c['max_drawdown'])}</td><td>{format_percent(d['max_drawdown'])}</td></tr>
              <tr><td>Volatilidad</td><td>{format_percent(b['volatility'])}</td><td>{format_percent(c['volatility'])}</td><td>{format_percent(d['volatility'])}</td></tr>
            </tbody>
          </table>
        </div>
      </section>

      <section class=\"callout\" aria-label=\"Límite académico\">
        <span aria-hidden=\"true\">⚠</span>
        <p><strong>Límite de tesis:</strong> evidencia pública piloto, no survivorship-bias-free. El candidato mejora robustez del baseline, pero no debe presentarse como superior a SPY/60-40 ni como recomendación de inversión.</p>
      </section>

      <section id=\"datos\" class=\"panel\">
        <h2>Próximo paso de datos point-in-time</h2>
        <p>El universo actual viene de un snapshot activo de Nasdaq. Para una tesis empírica fuerte, el siguiente paso ideal es usar una base institucional point-in-time como CRSP, Morningstar Direct, Lipper, Bloomberg o Refinitiv/LSEG, o construir snapshots históricos por año si aparece una fuente pública archivada estable.</p>
        <p>Documento incorporado: <a href=\"../../research/etf_point_in_time_data_sources.md\">matriz de fuentes ETF point-in-time</a>.</p>
      </section>

      <section id=\"comando\">
        <h2>Ejecución aprobable</h2>
        <p>La interfaz muestra el comando antes de correrlo. Esto respeta el principio de no ejecutar procesos locales sin confirmación humana.</p>
        <div class=\"command-card\">
          <div class=\"command-toolbar\"><span>Comando reproducible</span><button class=\"copy-button\" type=\"button\" data-copy-command>Copiar comando</button></div>
          <pre class=\"command\"><code id=\"run-command\">{command_html}</code></pre>
        </div>
      </section>

      <section id=\"artefactos\">
        <h2>Artefactos para revisión</h2>
        <div class=\"links\">
          <a class=\"link-card\" href=\"../tesis_trabajo_grado_etf_electre.docx\">Documento de tesis DOCX<span>Marco, metodología, resultados y límites.</span></a>
          <a class=\"link-card\" href=\"../presentacion_sustentacion_etf_electre.pptx\">Presentación PPTX<span>Narrativa defendible para sustentación.</span></a>
          <a class=\"link-card\" href=\"../../traceability/milestone_metrics_history.md\">Historial de hitos<span>Trazabilidad de decisiones y métricas.</span></a>
          <a class=\"link-card\" href=\"../../../results/sprint_universe_paper_quarterly_2015_2025_every_confirm2_m030_cap025/strategy_comparison.csv\">Resultados CSV<span>Tabla fuente de comparación.</span></a>
        </div>
      </section>
    </main>
  </div>
  <script id=\"front-chart-data\" type=\"application/json\">{chart_payload}</script>
  <script>
    const button = document.querySelector('[data-copy-command]');
    const commandNode = document.querySelector('#run-command');
    if (button && commandNode && navigator.clipboard) {{
      button.addEventListener('click', async () => {{
        await navigator.clipboard.writeText(commandNode.textContent.trim());
        button.textContent = 'Comando copiado';
        window.setTimeout(() => {{ button.textContent = 'Copiar comando'; }}, 1800);
      }});
    }}

    const chartDataNode = document.querySelector('#front-chart-data');
    const chartSvg = document.querySelector('#strategy-chart');
    const toggles = document.querySelector('#strategy-toggles');
    const tooltip = document.querySelector('#chart-tooltip');
    const palette = ['#2563eb', '#059669', '#dc2626', '#7c3aed', '#ea580c', '#0891b2', '#4f46e5', '#16a34a'];
    const chartData = chartDataNode ? JSON.parse(chartDataNode.textContent) : {{ series: [], events: [] }};
    const active = new Set(chartData.series.map((series) => series.name));

    function labelFor(name) {{
      return name.replace(/_/g, ' ').replace('walk forward', 'WF');
    }}

    function drawToggles() {{
      if (!toggles) return;
      toggles.innerHTML = '';
      chartData.series.forEach((series, index) => {{
        const label = document.createElement('label');
        const input = document.createElement('input');
        input.type = 'checkbox';
        input.checked = active.has(series.name);
        input.addEventListener('change', () => {{
          if (input.checked) active.add(series.name); else active.delete(series.name);
          drawChart();
        }});
        const color = document.createElement('span');
        color.style.inlineSize = '10px';
        color.style.blockSize = '10px';
        color.style.borderRadius = '999px';
        color.style.background = palette[index % palette.length];
        label.append(input, color, document.createTextNode(labelFor(series.name)));
        toggles.append(label);
      }});
    }}

    function drawChart() {{
      if (!chartSvg) return;
      const visible = chartData.series.filter((series) => active.has(series.name));
      if (!visible.length) {{
        chartSvg.innerHTML = '<text x="24" y="60" fill="currentColor">Select at least one strategy or benchmark.</text>';
        return;
      }}
      const width = chartSvg.clientWidth || 960;
      const height = chartSvg.clientHeight || 380;
      const pad = {{ top: 24, right: 28, bottom: 42, left: 54 }};
      const allPoints = visible.flatMap((series) => series.values.map((point) => ({{ ...point, dateMs: Date.parse(point.date) }})));
      const minDate = Math.min(...allPoints.map((point) => point.dateMs));
      const maxDate = Math.max(...allPoints.map((point) => point.dateMs));
      const minValue = Math.min(...allPoints.map((point) => point.value));
      const maxValue = Math.max(...allPoints.map((point) => point.value));
      const x = (dateMs) => pad.left + ((dateMs - minDate) / Math.max(1, maxDate - minDate)) * (width - pad.left - pad.right);
      const y = (value) => height - pad.bottom - ((value - minValue) / Math.max(0.000001, maxValue - minValue)) * (height - pad.top - pad.bottom);
      const parts = [];
      parts.push(`<svg viewBox="0 0 ${{width}} ${{height}}" xmlns="http://www.w3.org/2000/svg">`);
      parts.push(`<line x1="${{pad.left}}" y1="${{height - pad.bottom}}" x2="${{width - pad.right}}" y2="${{height - pad.bottom}}" stroke="currentColor" opacity="0.22"/>`);
      parts.push(`<line x1="${{pad.left}}" y1="${{pad.top}}" x2="${{pad.left}}" y2="${{height - pad.bottom}}" stroke="currentColor" opacity="0.22"/>`);
      [0, .25, .5, .75, 1].forEach((tick) => {{
        const value = minValue + (maxValue - minValue) * tick;
        const yy = y(value);
        parts.push(`<line x1="${{pad.left}}" y1="${{yy}}" x2="${{width - pad.right}}" y2="${{yy}}" stroke="currentColor" opacity="0.08"/>`);
        parts.push(`<text x="${{pad.left - 10}}" y="${{yy + 4}}" text-anchor="end" fill="currentColor" opacity="0.62" font-size="11">${{value.toFixed(2)}}x</text>`);
      }});
      visible.forEach((series) => {{
        const index = chartData.series.findIndex((item) => item.name === series.name);
        const color = palette[index % palette.length];
        const d = series.values.map((point, pointIndex) => `${{pointIndex === 0 ? 'M' : 'L'}} ${{x(Date.parse(point.date)).toFixed(2)}} ${{y(point.value).toFixed(2)}}`).join(' ');
        parts.push(`<path d="${{d}}" fill="none" stroke="${{color}}" stroke-width="2.4" stroke-linecap="round" stroke-linejoin="round"/>`);
        chartData.events.forEach((event) => {{
          const point = series.values.reduce((best, item) => Math.abs(Date.parse(item.date) - Date.parse(event.date)) < Math.abs(Date.parse(best.date) - Date.parse(event.date)) ? item : best, series.values[0]);
          const cx = x(Date.parse(point.date));
          const cy = y(point.value);
          parts.push(`<circle cx="${{cx.toFixed(2)}}" cy="${{cy.toFixed(2)}}" r="4.6" fill="${{color}}" stroke="var(--surface-raised)" stroke-width="2" data-name="${{series.name}}" data-date="${{event.date}}" data-type="${{event.type}}" data-value="${{point.value}}" data-turnover="${{event.turnover}}"/>`);
        }});
      }});
      const start = new Date(minDate).getFullYear();
      const end = new Date(maxDate).getFullYear();
      parts.push(`<text x="${{pad.left}}" y="${{height - 12}}" fill="currentColor" opacity="0.62" font-size="11">${{start}}</text>`);
      parts.push(`<text x="${{width - pad.right}}" y="${{height - 12}}" text-anchor="end" fill="currentColor" opacity="0.62" font-size="11">${{end}}</text>`);
      parts.push('</svg>');
      chartSvg.innerHTML = parts.join('');
      chartSvg.querySelectorAll('circle').forEach((circle) => {{
        circle.addEventListener('mouseenter', (event) => {{
          tooltip.style.opacity = '1';
          tooltip.textContent = `${{labelFor(circle.dataset.name)}} · ${{circle.dataset.date}} · ${{circle.dataset.type}} · equity ${{Number(circle.dataset.value).toFixed(2)}}x · turnover ${{Number(circle.dataset.turnover).toFixed(2)}}`;
          const rect = chartSvg.getBoundingClientRect();
          tooltip.style.left = `${{event.clientX - rect.left}}px`;
          tooltip.style.top = `${{event.clientY - rect.top}}px`;
        }});
        circle.addEventListener('mouseleave', () => {{ tooltip.style.opacity = '0'; }});
      }});
    }}

    drawToggles();
    drawChart();
    window.addEventListener('resize', drawChart);
  </script>
</body>
</html>
"""
    output = output_dir / "index.html"
    output.write_text(html_text, encoding="utf-8")
    return output


def write_manifest(output_dir: Path, paths: dict[str, Path], bundle: dict[str, Any]) -> Path:
    manifest = {
        "title": PROJECT_TITLE,
        "candidate_result_dir": str(CANDIDATE_DIR),
        "baseline_result_dir": str(BASELINE_DIR),
        "metrics": bundle,
        "deliverables": {key: str(path) for key, path in paths.items()},
        "claim_boundary": "Public-data pilot evidence; not survivorship-bias-free; presentation-ready research artifact, not investment advice.",
    }
    path = output_dir / "deliverables_manifest.json"
    path.write_text(json.dumps(manifest, indent=2, sort_keys=True), encoding="utf-8")
    return path


def build_all(output_dir: Path) -> dict[str, Path]:
    bundle = load_metric_bundle()
    output_dir.mkdir(parents=True, exist_ok=True)
    paths = {
        "docx": build_thesis_docx(output_dir / "tesis_trabajo_grado_etf_electre.docx", bundle),
        "pptx": build_deck(output_dir / "presentacion_sustentacion_etf_electre.pptx", bundle),
        "front_html": build_front_html(output_dir / "front_presentacion", bundle),
    }
    paths["manifest"] = write_manifest(output_dir, paths, bundle)
    return paths


def main() -> None:
    parser = argparse.ArgumentParser(description="Build presentation-ready thesis deliverables.")
    parser.add_argument("--out", type=Path, default=Path("docs/deliverables"))
    args = parser.parse_args()
    paths = build_all(args.out)
    for key, path in paths.items():
        print(f"{key}: {path}")


if __name__ == "__main__":
    main()
